#!/usr/bin/env python3
"""Build AgentPV defence deck (.pptx) from ppt制作指南.md.

Usage (repo root):
    pip install python-pptx
    python scripts/render_presentation.py
    python scripts/render_presentation.py --out reports/AgentPV_Final_Presentation.pptx
    python scripts/render_presentation.py --verify
"""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_GUIDE = ROOT / "docs" / "ppt制作指南.md"
DEFAULT_OUT = ROOT / "reports" / "AgentPV_Final_Presentation.pptx"

NAVY = RGBColor(0x1A, 0x36, 0x5D)
BODY = RGBColor(0x2D, 0x37, 0x48)
GRAY = RGBColor(0x71, 0x85, 0x9A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
ZEBRA = RGBColor(0xED, 0xF2, 0xF7)
ACCENT = RGBColor(0x2B, 0x6C, 0xB0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.5)
CONTENT_TOP = Inches(1.25)
CONTENT_BOTTOM = Inches(7.0)
GAP = Inches(0.12)


@dataclass
class TextBlock:
    text: str


@dataclass
class TableBlock:
    headers: list[str]
    rows: list[list[str]]


ContentBlock = Union[TextBlock, TableBlock]


@dataclass
class SlideSpec:
    slide_id: str
    title: str
    body: str
    blocks: list[ContentBlock] = field(default_factory=list)
    figure_paths: list[Path] = field(default_factory=list)
    notes_cn: str = ""
    notes_en: str = ""

    @property
    def table_count(self) -> int:
        return sum(1 for b in self.blocks if isinstance(b, TableBlock))


def _split_cells(line: str) -> list[str]:
    line = line.strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|"):
        line = line[:-1]
    return [c.strip() for c in line.split("|")]


def _is_separator_row(cells: list[str]) -> bool:
    return bool(cells) and all(re.fullmatch(r":?-{2,}:?", c.replace(" ", "")) for c in cells if c)


def parse_markdown_table(lines: list[str]) -> TableBlock:
    parsed: list[list[str]] = []
    for line in lines:
        cells = _split_cells(line)
        if not cells or _is_separator_row(cells):
            continue
        parsed.append(cells)
    if not parsed:
        return TableBlock([], [])
    headers = parsed[0]
    rows = parsed[1:]
    width = len(headers)
    rows = [r + [""] * (width - len(r)) if len(r) < width else r[:width] for r in rows]
    return TableBlock(headers=headers, rows=rows)


def parse_body_blocks(body: str) -> list[ContentBlock]:
    blocks: list[ContentBlock] = []
    lines = body.splitlines()
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if stripped.startswith("|"):
            table_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            table = parse_markdown_table(table_lines)
            if table.headers:
                blocks.append(table)
            continue
        text_lines: list[str] = []
        while i < len(lines) and not lines[i].strip().startswith("|"):
            text_lines.append(lines[i])
            i += 1
        text = "\n".join(text_lines).strip()
        if text:
            blocks.append(TextBlock(text))
    return blocks


def parse_guide(path: Path) -> list[SlideSpec]:
    text = path.read_text(encoding="utf-8")
    start = text.find("### Slide 1 ·")
    if start == -1:
        raise ValueError("Could not find slide section in guide")
    end = text.find("\n## 四、")
    if end == -1:
        end = text.find("\n## 五、")
    if end == -1:
        end = len(text)
    text = text[start:end]

    headers = list(
        re.finditer(r"^### ((?:Slide \d+|Appendix A\d+)) · (.+)$", text, re.MULTILINE)
    )
    specs: list[SlideSpec] = []
    for i, m in enumerate(headers):
        block_start = m.start()
        block_end = headers[i + 1].start() if i + 1 < len(headers) else len(text)
        block = text[block_start:block_end]
        slide_id = m.group(1)
        title = m.group(2).strip()

        body_m = re.search(
            r"\*\*On-slide text \(English only\)\*\*:\s*\n\n```\n(.*?)```",
            block,
            re.DOTALL,
        )
        body = body_m.group(1).strip() if body_m else ""

        raw_paths = re.findall(r"`(reports/[^`]+\.png)`", block)
        figure_paths: list[Path] = []
        for rp in raw_paths:
            p = ROOT / rp
            if p.is_file() and p not in figure_paths:
                figure_paths.append(p)

        cn_m = re.search(
            r"\*\*旁白（中文，详细）\*\*\s*\n(.*?)(?=\n\*\*Narration|\n\*\*Anticipated|\n---|\Z)",
            block,
            re.DOTALL,
        )
        en_m = re.search(
            r"\*\*Narration \(English, detailed\)\*\*\s*\n(.*?)(?=\n\*\*Anticipated|\n---|\Z)",
            block,
            re.DOTALL,
        )
        notes_cn = _clean_note(cn_m.group(1) if cn_m else "")
        notes_en = _clean_note(en_m.group(1) if en_m else "")

        specs.append(
            SlideSpec(
                slide_id=slide_id,
                title=title,
                body=body,
                blocks=parse_body_blocks(body),
                figure_paths=figure_paths,
                notes_cn=notes_cn,
                notes_en=notes_en,
            )
        )
    return specs


def _clean_note(s: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", s.strip())


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _looks_numeric(text: str) -> bool:
    t = text.strip().replace(",", "")
    if not t or t in {"—", "-", "✅", "⚠", "N/A"}:
        return False
    return bool(re.fullmatch(r"[\d.%+\-✅⚠]+(?:\s*ms)?(?:\s*\([^)]*\))?", t)) or bool(
        re.search(r"^\d", t)
    )


def _set_cell_text(cell, text: str, *, font_size: int, bold: bool = False, color: RGBColor = BODY, align=PP_ALIGN.LEFT) -> None:
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def _fill_table_cell(cell, text: str, *, header: bool, zebra: bool, font_size: int) -> None:
    if header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_cell_text(cell, text, font_size=font_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER if len(text) < 12 else PP_ALIGN.LEFT)
        return
    cell.fill.solid()
    cell.fill.fore_color.rgb = ZEBRA if zebra else WHITE
    align = PP_ALIGN.RIGHT if _looks_numeric(text) else PP_ALIGN.LEFT
    color = ACCENT if text.strip() in {"✅", "⚠"} else BODY
    _set_cell_text(cell, text, font_size=font_size, color=color, align=align)


def _table_font_size(cols: int, rows: int) -> int:
    if cols >= 8 or rows >= 12:
        return 9
    if cols >= 6 or rows >= 9:
        return 10
    if cols >= 5:
        return 11
    return 12


def _estimate_text_height(text: str, width_in: float, font_size: int) -> float:
    lines = max(1, len(text.splitlines()))
    chars_per_line = max(20, int(width_in * 72 / font_size * 0.55))
    wrapped = 0
    for line in text.splitlines():
        wrapped += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
    line_h = font_size / 72 * 1.35
    return wrapped * line_h + 0.08


def _estimate_table_height(rows: int, font_size: int) -> float:
    row_h = max(0.28, font_size / 72 * 1.55)
    return rows * row_h + 0.05


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 18,
    bold_first_line: bool = False,
    color: RGBColor = BODY,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.splitlines() or [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(3)
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            if bold_first_line and i == 0 and line.strip():
                run.font.bold = True
                run.font.color.rgb = NAVY
                run.font.size = Pt(font_size + 2)
    return box


def _add_markdown_table(
    slide,
    table: TableBlock,
    left,
    top,
    width,
) -> float:
    n_rows = len(table.rows) + 1
    n_cols = len(table.headers)
    font_size = _table_font_size(n_cols, n_rows)
    height = Inches(_estimate_table_height(n_rows, font_size))
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table

    col_weights = []
    for c in range(n_cols):
        sample = table.headers[c] + " ".join(row[c] for row in table.rows if c < len(row))
        col_weights.append(max(1, len(sample)))
    total = sum(col_weights) or 1
    for c in range(n_cols):
        tbl.columns[c].width = int(width * (col_weights[c] / total))

    for c, header in enumerate(table.headers):
        _fill_table_cell(tbl.cell(0, c), header, header=True, zebra=False, font_size=font_size)
    for r, row in enumerate(table.rows, start=1):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            _fill_table_cell(tbl.cell(r, c), val, header=False, zebra=r % 2 == 0, font_size=font_size)
    return height


def _add_header(slide, slide_id: str, title: str, page: int, total: int) -> None:
    tag = slide_id.replace("Slide ", "§").replace("Appendix ", "App ")
    _add_textbox(slide, MARGIN_L, Inches(0.25), Inches(4), Inches(0.35), tag, font_size=11, color=GRAY)
    _add_textbox(
        slide,
        Inches(11.5),
        Inches(0.25),
        Inches(1.3),
        Inches(0.35),
        f"{page} / {total}",
        font_size=11,
        color=GRAY,
        align=PP_ALIGN.RIGHT,
    )
    _add_textbox(slide, MARGIN_L, Inches(0.55), Inches(12.3), Inches(0.55), title, font_size=24, color=NAVY)


def _add_picture(slide, path: Path, left, top, width, height) -> None:
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def _layout_images(slide, paths: list[Path], text_width: float) -> None:
    if not paths:
        return
    n = len(paths)
    margin = MARGIN_L
    top_base = CONTENT_TOP

    if n == 1:
        img_w = Inches(5.8)
        img_h = Inches(5.0)
        left = MARGIN_L + Inches(text_width) + Inches(0.2)
        _add_picture(slide, paths[0], left, top_base, img_w, img_h)
        return

    if n == 2:
        img_w = Inches(5.5)
        img_h = Inches(2.6)
        left = MARGIN_L + Inches(text_width) + Inches(0.15)
        _add_picture(slide, paths[0], left, top_base, img_w, img_h)
        _add_picture(slide, paths[1], left, top_base + img_h + Inches(0.15), img_w, img_h)
        return

    row_top = Inches(4.35)
    usable_w = SLIDE_W - margin * 2
    gap = Inches(0.12)
    img_w = (usable_w - gap * (n - 1)) / n
    img_h = CONTENT_BOTTOM - row_top
    for i, p in enumerate(paths[:4]):
        left = margin + i * (img_w + gap)
        _add_picture(slide, p, left, row_top, img_w, img_h)


def _render_blocks(
    slide,
    blocks: list[ContentBlock],
    *,
    left,
    width,
    top,
    bottom,
) -> None:
    y = top
    width_in = width / 914400 * 96 / 72  # EMU to inches approx

    for block in blocks:
        if y >= bottom - Inches(0.1):
            break
        if isinstance(block, TextBlock):
            font_size = 16 if len(block.text) > 600 else 17
            est_h = Inches(_estimate_text_height(block.text, width_in, font_size))
            max_h = bottom - y
            box_h = min(est_h, max_h)
            _add_textbox(slide, left, y, width, box_h, block.text, font_size=font_size)
            y += box_h + GAP
        elif isinstance(block, TableBlock):
            n_rows = len(block.rows) + 1
            font_size = _table_font_size(len(block.headers), n_rows)
            est_h = Inches(_estimate_table_height(n_rows, font_size))
            max_h = bottom - y
            if est_h > max_h and n_rows > 2:
                font_size = max(8, font_size - 1)
                est_h = Inches(_estimate_table_height(n_rows, font_size))
            table_h = min(est_h, max_h)
            shape = slide.shapes.add_table(n_rows, len(block.headers), left, y, width, table_h)
            tbl = shape.table
            col_weights = [
                max(1, len(block.headers[c]) + sum(len(row[c]) for row in block.rows if c < len(row)))
                for c in range(len(block.headers))
            ]
            total = sum(col_weights) or 1
            for c in range(len(block.headers)):
                tbl.columns[c].width = int(width * (col_weights[c] / total))
            for c, header in enumerate(block.headers):
                _fill_table_cell(tbl.cell(0, c), header, header=True, zebra=False, font_size=font_size)
            for r, row in enumerate(block.rows, start=1):
                for c in range(len(block.headers)):
                    val = row[c] if c < len(row) else ""
                    _fill_table_cell(tbl.cell(r, c), val, header=False, zebra=r % 2 == 0, font_size=font_size)
            y += table_h + GAP


def _build_slide(prs: Presentation, spec: SlideSpec, page: int, total: int) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _set_slide_bg(slide, BG)

    if spec.slide_id == "Slide 1":
        body = spec.body.replace("[Your Name]", "Mansy")
        _add_textbox(
            slide,
            Inches(1.0),
            Inches(2.0),
            Inches(11.3),
            Inches(3.5),
            body,
            font_size=22,
            bold_first_line=True,
            align=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            Inches(11.5),
            Inches(0.25),
            Inches(1.3),
            Inches(0.35),
            f"{page} / {total}",
            font_size=11,
            color=GRAY,
            align=PP_ALIGN.RIGHT,
        )
    else:
        _add_header(slide, spec.slide_id, spec.title, page, total)
        has_imgs = bool(spec.figure_paths)
        bottom_imgs = has_imgs and len(spec.figure_paths) >= 3
        text_w = 6.0 if has_imgs and len(spec.figure_paths) <= 2 else 12.3
        content_bottom = Inches(4.25) if bottom_imgs else CONTENT_BOTTOM
        content_width = Inches(text_w)

        _render_blocks(
            slide,
            spec.blocks,
            left=MARGIN_L,
            width=content_width,
            top=CONTENT_TOP,
            bottom=content_bottom,
        )

        if has_imgs:
            if bottom_imgs:
                _layout_images(slide, spec.figure_paths, text_width=12.3)
            else:
                _layout_images(slide, spec.figure_paths, text_width=text_w)

    notes_parts = []
    if spec.notes_cn:
        notes_parts.append("【中文旁白】\n" + spec.notes_cn)
    if spec.notes_en:
        notes_parts.append("【English Narration】\n" + spec.notes_en)
    if notes_parts:
        slide.notes_slide.notes_text_frame.text = "\n\n".join(notes_parts)


def verify_presentation(guide: Path, pptx: Path) -> list[str]:
    specs = parse_guide(guide)
    if not pptx.is_file():
        return [f"PPTX not found: {pptx}"]

    prs = Presentation(str(pptx))
    issues: list[str] = []

    if len(prs.slides) != len(specs):
        issues.append(f"Slide count mismatch: guide={len(specs)} pptx={len(prs.slides)}")

    for i, (spec, slide) in enumerate(zip(specs, prs.slides), start=1):
        tables_on_slide = sum(1 for sh in slide.shapes if sh.has_table)
        if spec.table_count != tables_on_slide:
            issues.append(
                f"{spec.slide_id} ({spec.title}): expected {spec.table_count} table(s), found {tables_on_slide}"
            )
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text
            bad_lines = [
                ln
                for ln in txt.splitlines()
                if ln.strip().startswith("|") and ln.count("|") >= 2
            ]
            if bad_lines:
                preview = bad_lines[0][:80].encode("ascii", "replace").decode()
                issues.append(
                    f"{spec.slide_id}: raw markdown table still in text box — '{preview}...'"
                )

    missing_imgs = []
    for spec in specs:
        for p in spec.figure_paths:
            if not p.is_file():
                missing_imgs.append(f"{spec.slide_id}: missing figure {p}")
    issues.extend(missing_imgs)

    return issues


def build_presentation(guide: Path, out: Path) -> int:
    specs = parse_guide(guide)
    if not specs:
        raise SystemExit(f"No slides parsed from {guide}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = len(specs)

    for i, spec in enumerate(specs, start=1):
        _build_slide(prs, spec, i, total)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    with_figs = sum(1 for s in specs if s.figure_paths)
    with_tables = sum(1 for s in specs if s.table_count)
    print(
        f"[ok] Wrote {out} ({total} slides, {with_figs} with PNGs, {with_tables} with tables)"
    )
    return total


def main() -> None:
    ap = argparse.ArgumentParser(description="Render AgentPV defence PPTX from ppt制作指南.md")
    ap.add_argument("--guide", type=Path, default=DEFAULT_GUIDE)
    ap.add_argument("--out", type=Path, default=DEFAULT_OUT)
    ap.add_argument(
        "--verify",
        action="store_true",
        help="After build (or on existing --out), check tables and slide count",
    )
    ap.add_argument("--verify-only", action="store_true", help="Verify existing PPTX without rebuilding")
    args = ap.parse_args()
    if not args.guide.is_file():
        raise SystemExit(f"Guide not found: {args.guide}")

    if not args.verify_only:
        build_presentation(args.guide, args.out)

    if args.verify or args.verify_only:
        issues = verify_presentation(args.guide, args.out)
        if issues:
            print("[verify] FAILED:")
            for item in issues:
                safe = item.encode("ascii", "replace").decode()
                print(f"  - {safe}")
            raise SystemExit(1)
        print(f"[verify] OK — {args.out}")


if __name__ == "__main__":
    main()
